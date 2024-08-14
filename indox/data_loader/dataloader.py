class DataLoader:
    """
    A class to load various types of files based on their extension.

    Attributes
    ----------
    file_path : str
        Path to the file to be loaded.
    extension : str
        File extension of the file to be loaded.

    Methods
    -------
    load_file()
        Loads the file based on its extension.
    """

    def __init__(self, file_path):
        """
        Parameters
        ----------
        file_path : str
            Path to the file to be loaded.
        """
        self.file_path = file_path
        self.extension = file_path.split('.')[-1].lower()

    def load_file(self):
        """
        Loads the file based on its extension.

        Returns
        -------
        object
            The content of the file, type depends on the file format.

        Raises
        ------
        ValueError
            If the file extension is not supported.
        """
        loaders = {
            'pdf': self._load_pdf,
            'xlsx': self._load_excel,
            'docx': self._load_docx,
            'csv': self._load_csv,
            'pt': self._load_pt,
            'txt': self._load_txt,
            'json': self._load_json,
            'xml': self._load_xml,
            'html': self._load_html,
            'md': self._load_md,
            'rtf': self._load_rtf,
            'pptx': self._load_pptx,
            'odp': self._load_odp,
            'odt': self._load_odt,
            'ods': self._load_ods,
            'epub': self._load_epub,
            'yml': self._load_yaml,
            'yaml': self._load_yaml,
            'sql': self._load_sql,
            'db': self._load_db,
            'sqlite': self._load_sqlite,
            'tar': self._load_tar,
            'zip': self._load_zip,
            'rar': self._load_rar,
            'png': self._load_image,
            'jpg': self._load_image,
            'jpeg': self._load_image,
            'tiff': self._load_image,
            'tif': self._load_image,
            'jsonl': self._load_jsonl,
            'ipynb': self._load_ipynb,
            'pkl': self._load_pkl,
            'npy': self._load_npy,
            'mat': self._load_mat,
        }

        if self.extension in loaders:
            return loaders[self.extension]()
        else:
            raise ValueError(f"Unsupported file extension: {self.extension}")

    def _load_pdf(self):
        import PyPDF2
        try:
            with open(self.file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                text = ''
                for page in reader.pages:
                    text += page.extract_text()
            return text
        except Exception as e:
            raise RuntimeError(f"Error loading PDF file: {e}")

    def _load_excel(self):
        import pandas as pd
        try:
            return pd.read_excel(self.file_path)
        except Exception as e:
            raise RuntimeError(f"Error loading Excel file: {e}")

    def _load_docx(self):
        from docx import Document
        try:
            doc = Document(self.file_path)
            return '\n'.join([p.text for p in doc.paragraphs])
        except Exception as e:
            raise RuntimeError(f"Error loading DOCX file: {e}")

    def _load_csv(self):
        import pandas as pd
        try:
            return pd.read_csv(self.file_path)
        except Exception as e:
            raise RuntimeError(f"Error loading CSV file: {e}")

    def _load_pt(self):
        import torch
        try:
            return torch.load(self.file_path)
        except Exception as e:
            raise RuntimeError(f"Error loading PyTorch tensor file: {e}")

    def _load_txt(self):
        try:
            with open(self.file_path, 'r') as f:
                return f.read()
        except Exception as e:
            raise RuntimeError(f"Error loading text file: {e}")

    def _load_json(self):
        import json
        try:
            with open(self.file_path, 'r') as f:
                return json.load(f)
        except Exception as e:
            raise RuntimeError(f"Error loading JSON file: {e}")

    def _load_xml(self):
        import xml.etree.ElementTree as ET
        try:
            tree = ET.parse(self.file_path)
            return tree.getroot()
        except Exception as e:
            raise RuntimeError(f"Error loading XML file: {e}")

    def _load_html(self):
        try:
            with open(self.file_path, 'r') as f:
                return f.read()
        except Exception as e:
            raise RuntimeError(f"Error loading HTML file: {e}")

    def _load_md(self):
        import markdown
        from bs4 import BeautifulSoup
        try:
            with open(self.file_path, 'r') as f:
                text = f.read()
            html_content = markdown.markdown(text)
            soup = BeautifulSoup(html_content, 'html.parser')
            text_content = soup.get_text()
            code_blocks = [f"```python\n{code.get_text()}\n```" for code in soup.find_all('code')]
            return text_content + "\n\n" + "\n\n".join(code_blocks)
        except Exception as e:
            raise RuntimeError(f"Error loading Markdown file: {e}")

    def _load_rtf(self):
        try:
            with open(self.file_path, 'r') as f:
                return f.read()
        except Exception as e:
            raise RuntimeError(f"Error loading RTF file: {e}")

    def _load_pptx(self):
        from pptx import Presentation
        try:
            ppt = Presentation(self.file_path)
            text = []
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return '\n'.join(text)
        except Exception as e:
            raise RuntimeError(f"Error loading PowerPoint file: {e}")

    def _load_odp(self):
        from pptx import Presentation
        try:
            ppt = Presentation(self.file_path)
            text = []
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return '\n'.join(text)
        except Exception as e:
            raise RuntimeError(f"Error loading ODP file: {e}")

    def _load_odt(self):
        import odf.opendocument
        from odf.text import P
        try:
            textdoc = odf.opendocument.load(self.file_path)
            text_elements = textdoc.getElementsByType(P)
            return '\n'.join([str(p) for p in text_elements])
        except Exception as e:
            raise RuntimeError(f"Error loading ODT file: {e}")

    def _load_ods(self):
        import pandas as pd
        try:
            return pd.read_excel(self.file_path, engine='odf')
        except Exception as e:
            raise RuntimeError(f"Error loading ODS file: {e}")

    def _load_epub(self):

        from ebooklib import epub
        from bs4 import BeautifulSoup

        try:
            book = epub.read_epub(self.file_path)
            text = []

            for item in book.get_items():
                if item.get_type() == epub.EpubHtml:
                    soup = BeautifulSoup(item.get_content(), 'html.parser')
                    text.append(soup.get_text())

            return '\n'.join(text)

        except Exception as e:
            raise RuntimeError(f"Error loading EPUB file: {e}")
    def _load_yaml(self):
        import yaml
        try:
            with open(self.file_path, 'r') as f:
                return yaml.safe_load(f)
        except Exception as e:
            raise RuntimeError(f"Error loading YAML file: {e}")

    def _load_sql(self):
        try:
            with open(self.file_path, 'r') as f:
                return f.read()
        except Exception as e:
            raise RuntimeError(f"Error loading SQL file: {e}")

    def _load_db(self):

        import sqlite3

        try:
            conn = sqlite3.connect(self.file_path)
            cursor = conn.cursor()
            cursor.execute("SELECT name FROM sqlite_master WHERE type='table';")
            tables = cursor.fetchall()
            conn.close()
            return tables
        except Exception as e:
            raise RuntimeError(f"Error loading SQLite database: {e}")

    def _load_sqlite(self):

        return self._load_db()
    def _load_tar(self):
        import tarfile
        try:
            with tarfile.open(self.file_path, 'r') as tar:
                tar.extractall()
            return f"Extracted {self.file_path}"
        except Exception as e:
            raise RuntimeError(f"Error extracting TAR file: {e}")

    def _load_zip(self):
        import zipfile
        try:
            with zipfile.ZipFile(self.file_path, 'r') as zip_ref:
                zip_ref.extractall()
            return f"Extracted {self.file_path}"
        except Exception as e:
            raise RuntimeError(f"Error extracting ZIP file: {e}")

    def _load_rar(self):
        import rarfile
        try:
            with rarfile.RarFile(self.file_path, 'r') as rar_ref:
                rar_ref.extractall()
            return f"Extracted {self.file_path}"
        except Exception as e:
            raise RuntimeError(f"Error extracting RAR file: {e}")

    def _load_image(self):
        from PIL import Image
        try:
            return Image.open(self.file_path)
        except Exception as e:
            raise RuntimeError(f"Error loading image file: {e}")

    def _load_jsonl(self):
        import json
        try:
            with open(self.file_path, 'r') as f:
                return [json.loads(line) for line in f]
        except Exception as e:
            raise RuntimeError(f"Error loading JSONL file: {e}")

    def _load_ipynb(self):
        import json
        try:
            with open(self.file_path, 'r') as f:
                return json.load(f)
        except Exception as e:
            raise RuntimeError(f"Error loading Jupyter Notebook file: {e}")

    def _load_pkl(self):
        import pickle
        try:
            with open(self.file_path, 'rb') as f:
                return pickle.load(f)
        except Exception as e:
            raise RuntimeError(f"Error loading pickle file: {e}")

    def _load_npy(self):
        import numpy as np
        try:
            return np.load(self.file_path, allow_pickle=True)
        except Exception as e:
            raise RuntimeError(f"Error loading NPY file: {e}")

    def _load_mat(self):
        from scipy.io import loadmat
        try:
            return loadmat(self.file_path)
        except Exception as e:
            raise RuntimeError(f"Error loading MAT file: {e}")

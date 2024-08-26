from indox.core.document_object import Document
from typing import List, Dict
import os


class Pptx:
    """
    Load a PowerPoint (.pptx) file and extract its text and metadata.

    Parameters:
    - file_path (str): The path to the PowerPoint file to be loaded.

    Methods:
    - load_file(): Loads the PowerPoint file and returns a list of `Document` objects, each containing
                   the text content of a slide and the associated metadata.

    Raises:
    - RuntimeError: If there is an error in loading the PowerPoint file.

    Notes:
    - Each slide's text is extracted and stored in a separate `Document` object.
    - Metadata includes the file name and the number of slides in the presentation.
    """

    def __init__(self, file_path: str):
        self.file_path = os.path.abspath(file_path)

    def load(self) -> List[Document]:
        from pptx import Presentation
        try:
            presentation = Presentation(self.file_path)
            documents = []

            # Create metadata
            metadata_dict = {
                'source': self.file_path,
                'num_slides': len(presentation.slides),
            }

            for i, slide in enumerate(presentation.slides):
                text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)

                # Create a Document for each slide
                document = Document(page_content='\n'.join(text), slide_number=i + 1, **metadata_dict)
                documents.append(document)

            return documents
        except Exception as e:
            raise RuntimeError(f"Error loading PowerPoint file: {e}")

    def load_and_split(self, splitter, remove_stopwords=False):
        from indox.data_loader.utils import load_and_process_input
        return load_and_process_input(loader=self.load, splitter=splitter, remove_stopwords=remove_stopwords)

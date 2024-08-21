import os
import io
import pandas as pd
import nbformat
import PyPDF2
from pptx import Presentation
from docx import Document
import json
import xml.etree.ElementTree as ET
from bs4 import BeautifulSoup
from google.oauth2.credentials import Credentials
from google_auth_oauthlib.flow import InstalledAppFlow
from google.auth.transport.requests import Request
from googleapiclient.discovery import build

class GoogleDrive:
    SCOPES = ['https://www.googleapis.com/auth/drive.readonly']

    def __init__(self, creds_file='tokenGoogleDrive.json', credentials_json='credentials.json'):
        """
        Initialize the GoogleDrive object and authenticate using OAuth2.

        Parameters:
        - creds_file (str): The path to the file where the credentials are stored.
        - credentials_json (str): The path to the JSON file containing the client secrets.

        Returns:
        - None
        """
        self.creds_file = creds_file
        self.credentials_json = credentials_json
        self.creds = self._authenticate()

    def _authenticate(self):
        """
        Authenticate the user and refresh credentials if necessary.

        Returns:
        - Credentials: The authenticated Google API credentials.
        """
        creds = None
        if os.path.exists(self.creds_file):
            print("Loading credentials from file.")
            creds = Credentials.from_authorized_user_file(self.creds_file, self.SCOPES)
        if not creds or not creds.valid:
            if creds and creds.expired and creds.refresh_token:
                print("Refreshing expired credentials.")
                try:
                    creds.refresh(Request())
                except Exception as e:
                    print(f"Failed to refresh credentials: {e}")
                    raise
            else:
                print("Authenticating new credentials.")
                try:
                    flow = InstalledAppFlow.from_client_secrets_file(self.credentials_json, self.SCOPES)
                    creds = flow.run_local_server(port=0)
                except Exception as e:
                    print(f"Failed to authenticate: {e}")
                    raise
            try:
                with open(self.creds_file, 'w') as token:
                    token.write(creds.to_json())
                print("Credentials saved to token file.")
            except Exception as e:
                print(f"Failed to save credentials: {e}")
                raise
        return creds

    def read(self, file_id):
        """
        Read and print content from a file in Google Drive based on its MIME type.

        Parameters:
        - file_id (str): The ID of the file in Google Drive to be read.

        Returns:
        - None
        """
        try:
            service = build('drive', 'v3', credentials=self.creds)
            file = service.files().get(fileId=file_id, fields='id, name, mimeType, webViewLink').execute()

            file_name = file.get('name')
            file_mime_type = file.get('mimeType')

            print(f"File Name: {file_name}")
            print(f"MIME Type: {file_mime_type}")
            print(f"Web View Link: {file.get('webViewLink')}")

            request = service.files().get_media(fileId=file_id)
            file_content = request.execute()

            if file_mime_type == 'text/plain':  # .txt, .py
                print("File Content:\n")
                print(file_content.decode('utf-8'))

            elif file_name.endswith('.ipynb'):  # .ipynb
                nb = nbformat.reads(file_content.decode('utf-8'), as_version=4)
                print("Jupyter Notebook Content:\n")
                for cell in nb.cells:
                    if cell.cell_type == 'code':
                        print(f"Code:\n{cell.source}\n")

            elif file_mime_type == 'application/pdf':  # .pdf
                pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_content))
                print("PDF Content:\n")
                for page in pdf_reader.pages:
                    print(page.extract_text() or "No text found on this page.")

            elif file_mime_type == 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet':  # .xlsx
                df = pd.read_excel(io.BytesIO(file_content))
                print("Excel Content:\n")
                print(df)

            elif file_mime_type == 'text/csv':  # .csv
                df = pd.read_csv(io.BytesIO(file_content))
                print("CSV Content:\n")
                print(df)

            elif file_mime_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':  # PowerPoint
                ppt = Presentation(io.BytesIO(file_content))
                print("PowerPoint Content:\n")
                for slide in ppt.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            print(shape.text)

            elif file_mime_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':  # .docx
                doc = Document(io.BytesIO(file_content))
                print("Word Document Content:\n")
                for paragraph in doc.paragraphs:
                    print(paragraph.text)

            elif file_mime_type == 'application/json':  # .json
                json_content = json.loads(file_content.decode('utf-8'))
                print("JSON Content:\n")
                print(json.dumps(json_content, indent=4))

            elif file_mime_type == 'application/xml':  # .xml
                tree = ET.ElementTree(ET.fromstring(file_content.decode('utf-8')))
                root = tree.getroot()
                print("XML Content:\n")
                ET.dump(root)

            elif file_mime_type == 'text/html':  # .html
                soup = BeautifulSoup(file_content, 'html.parser')
                print("HTML Content:\n")
                print(soup.prettify())

            else:
                print(f"Unsupported MIME type: {file_mime_type}")

        except Exception as err:
            print(f"An error occurred: {err}")
        finally:
            if os.path.exists(self.creds_file):
                os.remove(self.creds_file)



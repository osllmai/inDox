from indoxArcg.core.document_object import Document
from typing import List
import os


class Pptx:
    """
    Load a PowerPoint (.pptx) file and extract its text and metadata.

    Parameters:
    - include_notes (bool): Whether to include notes from each slide (default: False)

    Methods:
    - load(file_path): Loads the PowerPoint file and returns a list of `Document` objects, each containing
                       the text content of a slide and the associated metadata.
    """

    def __init__(self, include_notes: bool = False):
        self.include_notes = include_notes

    def load(self, file_path: str) -> List[Document]:
        from pptx import Presentation

        try:
            presentation = Presentation(file_path)
            documents = []

            metadata_dict = {
                "source": os.path.abspath(file_path),
                "num_slides": len(presentation.slides),
            }

            for i, slide in enumerate(presentation.slides):
                text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)

                if self.include_notes and slide.has_notes_slide:
                    notes_slide = slide.notes_slide
                    if notes_slide and notes_slide.notes_text_frame:
                        notes_text = notes_slide.notes_text_frame.text
                        text.append("\n--- Notes ---\n" + notes_text)

                document = Document(
                    page_content="\n".join(text), slide_number=i + 1, **metadata_dict
                )
                documents.append(document)

            return documents
        except Exception as e:
            raise RuntimeError(f"Error loading PowerPoint file: {e}")

    def load_and_split(self, splitter, remove_stopwords=False):
        from indoxRag.data_loader.utils import load_and_process_input

        return load_and_process_input(
            loader=lambda: self.load(self.file_path),
            splitter=splitter,
            remove_stopwords=remove_stopwords,
        )
